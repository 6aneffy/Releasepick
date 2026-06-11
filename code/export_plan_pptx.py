"""Build PPTX from finalized CardNewsPlan (python-pptx)."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from models import CardNewsPlan

MALGUN = Path(r"C:\Windows\Fonts\malgun.ttf")


def _blank_layout(prs: Presentation):
    for lay in prs.slide_layouts:
        if lay.name and "blank" in lay.name.lower():
            return lay
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[-1]


def _set_ko_font(run, size_pt: int) -> None:
    run.font.size = Pt(size_pt)
    if MALGUN.exists():
        run.font.name = "Malgun Gothic"
    else:
        run.font.name = "Arial Unicode MS"


def build_plan_pptx_bytes(plan: CardNewsPlan) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    layout = _blank_layout(prs)

    # Title slide — series + head copy
    slide = prs.slides.add_slide(layout)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = plan.series_title
    _set_ko_font(p.runs[0], 40)
    p2 = tf.add_paragraph()
    p2.text = plan.head_copy
    _set_ko_font(p2.runs[0], 24)

    for s in plan.slides:
        slide = prs.slides.add_slide(layout)
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(6.5))
        tf = tb.text_frame
        tf.word_wrap = True
        head = tf.paragraphs[0]
        head.text = f"[{s.role}] {s.title}"
        _set_ko_font(head.runs[0], 28)
        for b in s.bullets:
            para = tf.add_paragraph()
            para.text = f"· {b}"
            para.level = 0
            _set_ko_font(para.runs[0], 18)
        if s.footnote:
            fn = tf.add_paragraph()
            fn.text = f"각주: {s.footnote}"
            _set_ko_font(fn.runs[0], 12)

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()
